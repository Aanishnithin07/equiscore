"""
EquiScore File Extractor — PDF & PPTX Content Extraction
==========================================================
Extracts text content from pitch deck files using:
  - PyMuPDF (fitz) for PDF files
  - python-pptx for PowerPoint PPTX files

The extractor preserves slide ordering, inserts slide-number markers
(e.g., "[SLIDE 3]"), captures speaker notes, detects image-only slides,
and validates minimum content thresholds before passing to the LLM.
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass, field

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ══════════════════════════════════════════════════════════════════════════
# Data Structures
# ══════════════════════════════════════════════════════════════════════════

@dataclass
class ExtractedContent:
    """
    Container for all content extracted from a pitch deck file.

    Attributes:
        raw_text: Concatenated text with slide markers like "[SLIDE 1]".
        slide_count: Total number of slides/pages in the document.
        has_images: Whether any images were detected in the document.
        image_count: Total number of images across all slides.
        word_count: Total word count of extracted text (excluding markers).
        extraction_method: Which library was used ("pymupdf" or "python-pptx").
        extraction_warnings: Issues detected during extraction (e.g., image-only slides).
    """

    raw_text: str
    slide_count: int
    has_images: bool
    image_count: int
    word_count: int
    extraction_method: str
    extraction_warnings: list[str] = field(default_factory=list)


class ExtractionError(Exception):
    """
    Raised when content extraction fails validation.

    Common causes:
      - File is corrupt or unreadable
      - Too few slides (< 3)
      - Insufficient text content (< 50 words — likely image-only deck)
    """

    pass


# ══════════════════════════════════════════════════════════════════════════
# PitchDeckExtractor
# ══════════════════════════════════════════════════════════════════════════

class PitchDeckExtractor:
    """
    Extracts and validates pitch deck content from PDF and PPTX files.

    Usage:
        extractor = PitchDeckExtractor()
        content = await extractor.extract(file_bytes, "pitch.pdf")
    """

    # Minimum thresholds for valid pitch decks
    MIN_WORD_COUNT: int = 50
    MIN_SLIDE_COUNT: int = 3
    # Threshold below which a slide is considered "image-only"
    IMAGE_ONLY_TEXT_THRESHOLD: int = 20

    async def extract(
        self,
        file_bytes: bytes,
        filename: str,
    ) -> ExtractedContent:
        """
        Extract content from a pitch deck file.

        Auto-detects format by filename extension and delegates to the
        appropriate extraction method.

        Args:
            file_bytes: Raw bytes of the uploaded file.
            filename: Original filename (used to determine format).

        Returns:
            ExtractedContent with all text, metadata, and warnings.

        Raises:
            ExtractionError: If the file is corrupt, too short, or unreadable.
            ValueError: If the file format is unsupported.
        """
        lower_filename = filename.lower()

        if lower_filename.endswith(".pdf"):
            content = self._extract_pdf(file_bytes)
        elif lower_filename.endswith(".pptx"):
            content = self._extract_pptx(file_bytes)
        else:
            raise ValueError(
                f"Unsupported file format: {filename}. "
                "Only .pdf and .pptx files are accepted."
            )

        # ── Post-extraction validation ────────────────────────────────
        self._validate_content(content)

        return content

    def _extract_pdf(self, file_bytes: bytes) -> ExtractedContent:
        """
        Extract text from a PDF file using PyMuPDF (fitz).

        For each page:
          - Extracts text using get_text("text") for clean content
          - Uses get_text("dict") to detect images and analyze layout
          - Detects image-only pages (< 20 chars of text but has images)

        Args:
            file_bytes: Raw PDF file bytes.

        Returns:
            ExtractedContent populated with text, slide counts, and warnings.

        Raises:
            ExtractionError: If PyMuPDF cannot open or parse the file.
        """
        warnings: list[str] = []
        text_parts: list[str] = []
        total_images = 0

        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
        except Exception as e:
            raise ExtractionError(
                f"Failed to open PDF file: {str(e)}. "
                "The file may be corrupt or password-protected."
            )

        slide_count = len(doc)

        for page_num in range(slide_count):
            page = doc[page_num]

            # ── Text extraction ───────────────────────────────────────
            page_text = page.get_text("text").strip()
            text_parts.append(f"[SLIDE {page_num + 1}]")

            if page_text:
                # Sanitize: strip null bytes and control characters
                sanitized = self._sanitize_text(page_text)
                text_parts.append(sanitized)

            # ── Layout analysis for image detection ───────────────────
            try:
                page_dict = page.get_text("dict")
                page_images = 0
                for block in page_dict.get("blocks", []):
                    # Block type 1 = image block in PyMuPDF
                    if block.get("type") == 1:
                        page_images += 1
                total_images += page_images

                # ── Image-only slide detection ────────────────────────
                if len(page_text) < self.IMAGE_ONLY_TEXT_THRESHOLD and page_images > 0:
                    warnings.append(
                        f"Slide {page_num + 1} appears to be image-only "
                        f"with no extractable text ({page_images} image(s) detected)"
                    )
            except Exception:
                # Layout analysis is best-effort — don't fail the whole extraction
                pass

            text_parts.append("")  # Blank line between slides

        doc.close()

        raw_text = "\n".join(text_parts)
        word_count = self._count_words(raw_text)

        return ExtractedContent(
            raw_text=raw_text,
            slide_count=slide_count,
            has_images=total_images > 0,
            image_count=total_images,
            word_count=word_count,
            extraction_method="pymupdf",
            extraction_warnings=warnings,
        )

    def _extract_pptx(self, file_bytes: bytes) -> ExtractedContent:
        """
        Extract text from a PPTX file using python-pptx.

        Extracts content from:
          - TextFrame shapes (body text, titles)
          - Table cells
          - Speaker notes (prepended with "[SPEAKER NOTES]")
          - Image shape detection and counting

        Args:
            file_bytes: Raw PPTX file bytes.

        Returns:
            ExtractedContent populated with text, slide counts, and warnings.

        Raises:
            ExtractionError: If python-pptx cannot open or parse the file.
        """
        warnings: list[str] = []
        text_parts: list[str] = []
        total_images = 0

        try:
            prs = Presentation(io.BytesIO(file_bytes))
        except Exception as e:
            raise ExtractionError(
                f"Failed to open PPTX file: {str(e)}. "
                "The file may be corrupt or in an unsupported format."
            )

        slide_count = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides, start=1):
            text_parts.append(f"[SLIDE {slide_idx}]")
            slide_text_fragments: list[str] = []
            slide_images = 0

            for shape in slide.shapes:
                # ── Text frames (titles, body text, etc.) ─────────────
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_text_fragments.append(
                                self._sanitize_text(para_text)
                            )

                # ── Tables ────────────────────────────────────────────
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_texts.append(self._sanitize_text(cell_text))
                        if row_texts:
                            slide_text_fragments.append(" | ".join(row_texts))

                # ── Image detection ───────────────────────────────────
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_images += 1

            total_images += slide_images

            # ── Append slide text ─────────────────────────────────────
            if slide_text_fragments:
                text_parts.append("\n".join(slide_text_fragments))

            # ── Image-only slide warning ──────────────────────────────
            combined_text = " ".join(slide_text_fragments)
            if len(combined_text) < self.IMAGE_ONLY_TEXT_THRESHOLD and slide_images > 0:
                warnings.append(
                    f"Slide {slide_idx} appears to be image-only "
                    f"with no extractable text ({slide_images} image(s) detected)"
                )

            # ── Speaker notes ─────────────────────────────────────────
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    sanitized_notes = self._sanitize_text(notes_text)
                    text_parts.append(f"[SPEAKER NOTES] {sanitized_notes}")

            text_parts.append("")  # Blank line between slides

        raw_text = "\n".join(text_parts)
        word_count = self._count_words(raw_text)

        return ExtractedContent(
            raw_text=raw_text,
            slide_count=slide_count,
            has_images=total_images > 0,
            image_count=total_images,
            word_count=word_count,
            extraction_method="python-pptx",
            extraction_warnings=warnings,
        )

    def _validate_content(self, content: ExtractedContent) -> None:
        """
        Validate that extracted content meets minimum thresholds.

        Args:
            content: The extracted content to validate.

        Raises:
            ExtractionError: If content is below minimum thresholds.
        """
        if content.slide_count < self.MIN_SLIDE_COUNT:
            raise ExtractionError(
                f"Pitch deck has only {content.slide_count} slide(s) — "
                f"minimum {self.MIN_SLIDE_COUNT} slides required."
            )

        if content.word_count < self.MIN_WORD_COUNT:
            raise ExtractionError(
                f"Insufficient content extracted ({content.word_count} words). "
                f"Minimum {self.MIN_WORD_COUNT} words required. "
                "The submission may be image-only or corrupt."
            )

    @staticmethod
    def _sanitize_text(text: str) -> str:
        """
        Sanitize extracted text to prevent LLM prompt injection and encoding issues.

        Strips:
          - Null bytes (\\x00) that can break JSON serialization
          - Control characters (except newlines and tabs for formatting)
          - Excessive whitespace

        Args:
            text: Raw text to sanitize.

        Returns:
            Cleaned text safe for LLM prompt injection.
        """
        # Remove null bytes
        text = text.replace("\x00", "")
        # Remove control characters except \n, \r, \t
        text = re.sub(r"[\x01-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
        # Collapse multiple spaces into one
        text = re.sub(r" {2,}", " ", text)
        # Collapse more than 2 consecutive newlines
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    @staticmethod
    def _count_words(text: str) -> int:
        """
        Count words in text, excluding slide markers like "[SLIDE 1]".

        Args:
            text: Text with slide markers.

        Returns:
            Word count of actual content (not markers).
        """
        # Remove slide markers before counting
        cleaned = re.sub(r"\[SLIDE \d+\]", "", text)
        cleaned = re.sub(r"\[SPEAKER NOTES\]", "", cleaned)
        return len(cleaned.split())
